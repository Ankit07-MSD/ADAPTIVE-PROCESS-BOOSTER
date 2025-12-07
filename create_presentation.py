#!/usr/bin/env python3
"""
Script to create PowerPoint presentation for Adaptive Process Booster
Requires: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed!")
    print("Please install it using: pip install python-pptx")
    exit(1)

def create_presentation():
    """Create PowerPoint presentation with all slides."""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    primary_color = RGBColor(46, 134, 171)  # Blue
    secondary_color = RGBColor(6, 167, 125)  # Green
    accent_color = RGBColor(241, 143, 1)  # Orange
    dark_text = RGBColor(50, 50, 50)
    
    # SLIDE 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Adaptive Process Booster"
    subtitle.text = "Dynamic Resource Manager - Advanced Edition\n\nA Comprehensive GUI-Based Process Management System\nReal-time Process Monitoring & Priority Optimization\nBuilt with Python, Tkinter, and psutil"
    
    title_shape = title.text_frame.paragraphs[0]
    title_shape.font.size = Pt(44)
    title_shape.font.bold = True
    title_shape.font.color.rgb = primary_color
    
    # SLIDE 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What is Adaptive Process Booster?"
    content.text = """Definition:
• Advanced process management application
• Real-time monitoring and control of system processes
• Intelligent priority optimization based on resource usage

Key Capabilities:
• Real-time process monitoring (1-second refresh)
• Dynamic priority boosting
• System resource tracking
• Process termination and management
• Comprehensive logging and history

Platform Support:
• Windows (with Administrator privileges)
• Linux/Mac (with appropriate permissions)"""
    
    # SLIDE 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Why Do We Need Process Management?"
    content.text = """Challenges:

1. Resource Competition
   • Multiple processes competing for CPU and memory
   • Critical applications may run slowly
   • System performance degradation

2. Manual Management
   • No easy way to identify resource-heavy processes
   • Difficult to prioritize important applications
   • Lack of real-time monitoring tools

3. System Optimization
   • Need for automatic resource allocation
   • Dynamic adjustment based on usage patterns
   • Historical tracking for analysis

Impact:
• Reduced productivity
• System slowdowns
• Poor user experience
• Inefficient resource utilization"""
    
    # SLIDE 4: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Solution - Adaptive Process Booster"
    content.text = """Main Objectives:

1. Real-time Monitoring
   • Continuous process scanning
   • Live system resource tracking
   • Instant updates every second

2. Intelligent Priority Management
   • Automatic priority boosting based on usage
   • Manual priority control
   • Multiple priority levels

3. User-Friendly Interface
   • Modern tabbed GUI
   • Search, filter, and sort capabilities
   • Detailed process information

4. Comprehensive Management
   • Process termination
   • History logging
   • Data export functionality

Key Innovation:
Dynamic Score Calculation: Score = 0.6 × CPU% + 0.4 × RAM%"""
    
    # SLIDE 5: Features Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Feature-Rich Process Management System"
    content.text = """Core Features:

1. Process Monitoring
   • Real-time process list with PID, Name, CPU%, RAM%, Score
   • Status tracking (running, sleeping, zombie)
   • Automatic refresh every second

2. Priority Management
   • 6 priority levels (Realtime to Low)
   • Manual priority boosting
   • Automatic priority optimization

3. Search & Filter
   • Search by process name or PID
   • Filter by CPU, RAM, or Score thresholds
   • Sort by any column

4. Process Details
   • Comprehensive process information
   • Memory usage (RSS/VMS)
   • Thread count, creation time
   • Executable path and command line

5. System Monitoring
   • CPU usage tracking
   • Memory usage (percentage and GB)
   • Disk usage monitoring
   • Process count display

6. History & Logging
   • Action history with timestamps
   • Export logs to text file
   • Clear history functionality"""
    
    # SLIDE 6: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture & Design"
    content.text = """Architecture:

GUI Layer (Tkinter)
├── Process Monitor Tab
├── System Resources Tab
└── History & Logs Tab
         │
         │ Thread-Safe Communication (root.after())
         │
Background Monitoring Thread
├── Process Scanning (psutil)
├── System Stats Collection
├── Score Calculation
└── Auto-Boost Logic
         │
         │
Operating System
├── Process Management
├── Priority Control
└── Resource Allocation

Key Technologies:
• Python 3.6+: Core programming language
• Tkinter: GUI framework
• psutil: System and process utilities
• Threading: Background monitoring
• CSV/JSON: Data export formats"""
    
    # SLIDE 7: Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Components Breakdown"
    content.text = """1. Priority Management Module
   • Cross-platform priority mapping
   • Windows: Priority classes (Realtime, High, Normal, etc.)
   • Linux: Nice values (-10 to +10)
   • Dynamic priority assignment

2. Process Monitoring Engine
   • Continuous process scanning
   • Real-time data collection
   • Score calculation algorithm
   • Auto-boost decision making

3. GUI Framework
   • Tabbed interface (3 main tabs)
   • Treeview for process display
   • Progress bars for system stats
   • Context menus and dialogs

4. Data Management
   • Process history tracking
   • Action logging system
   • CSV export functionality
   • Settings persistence

5. Threading System
   • Background monitoring thread
   • Thread-safe GUI updates
   • Non-blocking operations
   • Efficient resource usage"""
    
    # SLIDE 8: Score Algorithm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Intelligent Priority Decision Making"
    content.text = """Score Calculation Formula:
Score = 0.6 × CPU_Usage% + 0.4 × RAM_Usage%

Why This Formula?
• CPU Weight (60%): CPU is often the bottleneck
• RAM Weight (40%): Memory usage is also critical
• Weighted Approach: Balances both factors

Example Calculation:

Process A: CPU = 80%, RAM = 60%
Score = 0.6 × 80 + 0.4 × 60 = 48 + 24 = 72
Result: High score → Auto-boost candidate

Process B: CPU = 30%, RAM = 20%
Score = 0.6 × 30 + 0.4 × 20 = 18 + 8 = 26
Result: Low score → No boost needed

Auto-Boost Logic:
IF (Auto-Boost Enabled) AND (Score > Threshold):
    THEN Boost Process Priority
    LOG Action to History

Benefits:
• Automatic resource optimization
• No manual intervention needed
• Configurable threshold
• Real-time adaptation"""
    
    # SLIDE 9: GUI Design
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "User Interface - Three Main Tabs"
    content.text = """Tab 1: Process Monitor
• Process List Table
  - Columns: PID, Name, CPU%, RAM%, Score, Status
  - Sortable columns
  - Color-coded indicators
• Toolbar
  - Search box (real-time filtering)
  - Sort dropdown
  - Filter dropdown (High CPU/RAM/Score)
• Action Buttons
  - Boost Process
  - View Details
  - Kill Process
• Status Bar
  - Current selection
  - Operation status

Tab 2: System Resources
• CPU Usage (percentage + progress bar)
• Memory Usage (percentage + GB + progress bar)
• Disk Usage (percentage + GB + progress bar)
• Running Process Count
• Real-time updates

Tab 3: History & Logs
• Timestamped action log
• Scrollable text widget
• Export functionality
• Clear history option

Menu Bar:
• File: Export CSV, Exit
• Tools: Settings, Clear History"""
    
    # SLIDE 10: Advanced Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Advanced Capabilities"
    content.text = """1. Search & Filter System
   • Search: Real-time text search by name or PID
   • Filters: All, High CPU (>50%), High RAM (>50%), High Score (>50)
   • Sorting: Click column headers or use dropdown

2. Process Details Window
   • Complete process information
   • Memory breakdown (RSS/VMS)
   • Thread information
   • Creation timestamp
   • Executable path
   • Command line arguments
   • Working directory

3. Context Menu (Right-Click)
   • Quick access to actions
   • View Details
   • Boost Process
   • Kill Process

4. Settings Configuration
   • Enable/Disable Auto-Boost
   • Configure threshold value
   • Select priority level
   • Persistent settings

5. Data Export
   • Export process list to CSV
   • Export logs to text file
   • File dialog for save location

6. Selection Preservation
   • Maintains selection during updates
   • Scroll position preservation
   • Smooth user experience"""
    
    # SLIDE 11: System Monitoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Real-Time System Resource Tracking"
    content.text = """CPU Monitoring:
• Real-time CPU usage percentage
• Visual progress bar indicator
• Historical tracking (60 data points)
• Per-process CPU breakdown

Memory Monitoring:
• Total system memory
• Used memory (GB)
• Available memory (GB)
• Memory usage percentage
• Visual progress bar

Disk Monitoring:
• Total disk space
• Used disk space (GB)
• Free disk space (GB)
• Disk usage percentage
• Visual progress bar

Process Statistics:
• Total running processes
• Process status breakdown
• Resource consumption summary

Update Frequency:
• 1-second refresh interval
• Real-time data collection
• Thread-safe updates
• Non-blocking operations

Benefits:
• System health awareness
• Resource bottleneck identification
• Performance optimization insights"""
    
    # SLIDE 12: Priority Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Multi-Level Priority Control"
    content.text = """Priority Levels (Windows):
1. Realtime - Highest priority (use with caution)
2. High - Elevated priority for important tasks
3. Above Normal - Slightly elevated
4. Normal - Default system priority
5. Below Normal - Reduced priority
6. Low - Lowest priority for background tasks

Priority Levels (Linux):
1. Very High (Nice: -10)
2. High (Nice: -5)
3. Above Normal (Nice: -2)
4. Normal (Nice: 0)
5. Below Normal (Nice: +5)
6. Low (Nice: +10)

Usage Scenarios:
• Manual Boost: User selects process and priority
• Auto-Boost: Automatic based on score threshold
• Priority Selection Dialog: User-friendly interface

Implementation:
• Cross-platform compatibility
• Error handling for access denied
• Confirmation for critical operations
• Success/failure feedback"""
    
    # SLIDE 13: Threading
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Efficient Multi-Threading Architecture"
    content.text = """Threading Model:

Main Thread (GUI)          Background Thread (Monitoring)
     │                              │
     ├─ User Interactions           ├─ Process Scanning
     ├─ GUI Updates                 ├─ System Stats
     ├─ Event Handling              ├─ Score Calculation
     └─ Window Management           └─ Auto-Boost Logic
              │                              │
              └────── root.after() ──────────┘
              (Thread-Safe Communication)

Why Threading?
• Non-Blocking: GUI remains responsive
• Real-Time: Continuous monitoring
• Efficiency: Parallel operations
• User Experience: Smooth interface

Thread Safety:
• root.after() for GUI updates
• No direct GUI access from background thread
• Thread-safe data structures
• Proper synchronization

Performance Optimizations:
• 1-second refresh interval (balanced)
• Efficient psutil usage
• Limited history (100 entries max)
• Selection preservation (reduces redraws)
• Filtered updates (only visible processes)

Memory Management:
• Bounded history (deque with maxlen)
• Treeview item cleanup
• No memory leaks
• Efficient data structures"""
    
    # SLIDE 14: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Practical Applications"
    content.text = """1. Performance Optimization
   • Boost important applications
   • Improve system responsiveness
   • Optimize resource allocation
   • Example: Boost video editing software during rendering

2. System Monitoring
   • Identify resource-heavy processes
   • Monitor system health
   • Track resource usage trends
   • Example: Find processes causing slowdowns

3. Process Management
   • Terminate unresponsive processes
   • Manage background tasks
   • Control system resources
   • Example: Kill frozen applications

4. Development & Testing
   • Monitor application performance
   • Analyze resource consumption
   • Export data for analysis
   • Example: Performance profiling

5. System Administration
   • Server process management
   • Resource allocation
   • System optimization
   • Example: Server monitoring and control

6. Educational Purpose
   • Learn process management
   • Understand system resources
   • Study threading concepts
   • Example: Teaching operating systems

Target Users:
• System administrators
• Developers
• Power users
• IT professionals
• Students"""
    
    # SLIDE 15: Technical Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Deep Dive"
    content.text = """Key Functions:

1. boost_process_priority(pid, priority_level)
   • Sets process priority
   • Cross-platform implementation
   • Error handling

2. monitor_processes(gui_app)
   • Background monitoring loop
   • Process scanning
   • System stats collection
   • Auto-boost execution

3. calculate_score(cpu, ram)
   • Weighted score calculation
   • Returns priority indicator

4. update_process_list(data)
   • Thread-safe GUI update
   • Filtering and sorting
   • Selection preservation

5. get_process_details(pid)
   • Comprehensive process info
   • Memory breakdown
   • System information

Data Structures:
• deque for history (bounded)
• Dictionary for priority mapping
• List for process data
• String variables for GUI state

Error Handling:
• Process not found
• Access denied
• Invalid input
• File I/O errors
• Thread exceptions

Dependencies:
• psutil: System and process utilities
• tkinter: GUI framework
• threading: Multi-threading
• csv: Data export
• datetime: Timestamps"""
    
    # SLIDE 16: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Why Choose Adaptive Process Booster?"
    content.text = """User Benefits:
✅ Easy to Use: Intuitive GUI interface
✅ Real-Time: Live process monitoring
✅ Comprehensive: Multiple features in one tool
✅ Efficient: Optimized performance
✅ Flexible: Customizable settings
✅ Cross-Platform: Works on Windows and Linux

Technical Advantages:
✅ Thread-Safe: Proper multi-threading
✅ Scalable: Handles many processes
✅ Robust: Error handling and validation
✅ Extensible: Easy to add features
✅ Well-Structured: Clean code architecture

Performance Benefits:
✅ Automatic Optimization: Auto-boost feature
✅ Resource Awareness: System monitoring
✅ Historical Tracking: Logs and history
✅ Data Export: CSV and text formats

Comparison with Alternatives:
• More features than Task Manager
• Better GUI than command-line tools
• Real-time updates
• Automatic optimization
• Comprehensive logging"""
    
    # SLIDE 17: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Potential Improvements & Extensions"
    content.text = """Planned Features:

1. Visual Graphs & Charts
   • CPU usage graphs over time
   • Memory usage trends
   • Process comparison charts
   • Real-time plotting

2. Process Grouping
   • Group related processes
   • Batch operations
   • Process tree view
   • Parent-child relationships

3. Scheduled Tasks
   • Schedule priority changes
   • Time-based boosting
   • Automated workflows
   • Cron-like functionality

4. Network Monitoring
   • Network usage per process
   • Connection tracking
   • Bandwidth monitoring
   • Network priority

5. Advanced Filtering
   • Custom filter rules
   • Saved filter presets
   • Complex queries
   • Regular expressions

6. Notifications
   • Alert system
   • Threshold notifications
   • Email/SMS alerts
   • Desktop notifications

7. Remote Monitoring
   • Network monitoring
   • Remote process control
   • Multi-system dashboard
   • Client-server architecture

8. Machine Learning
   • Predictive priority adjustment
   • Usage pattern recognition
   • Intelligent auto-boost
   • Anomaly detection"""
    
    # SLIDE 18: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Summary & Key Takeaways"
    content.text = """Project Summary:
• Advanced process management system
• Real-time monitoring and control
• Intelligent priority optimization
• User-friendly GUI interface
• Comprehensive feature set

Key Achievements:
✅ Real-time process monitoring
✅ Dynamic priority management
✅ System resource tracking
✅ Advanced filtering and search
✅ Comprehensive logging
✅ Data export capabilities
✅ Cross-platform compatibility

Technical Highlights:
• Multi-threaded architecture
• Thread-safe GUI updates
• Efficient resource usage
• Robust error handling
• Clean code structure

Impact:
• Improved system performance
• Better resource utilization
• Enhanced user productivity
• Educational value
• Practical utility

Final Message:
"Adaptive Process Booster - Empowering users with intelligent process management and real-time system optimization." """
    
    # Save presentation
    filename = "Adaptive_Process_Booster_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    try:
        filename = create_presentation()
        print("\n🎉 PowerPoint presentation is ready!")
        print(f"📁 File location: {filename}")
        print("\n💡 Tips:")
        print("   - Open the file in Microsoft PowerPoint or LibreOffice Impress")
        print("   - Add screenshots of the application for better visuals")
        print("   - Customize colors and fonts as needed")
        print("   - Add diagrams and flowcharts for architecture slides")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\nTroubleshooting:")
        print("1. Make sure python-pptx is installed: pip install python-pptx")
        print("2. Check if you have write permissions in the current directory")
        print("3. Ensure no other program is using the output file")




